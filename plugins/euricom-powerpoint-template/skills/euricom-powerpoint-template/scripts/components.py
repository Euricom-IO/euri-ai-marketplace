#!/usr/bin/env python3
"""
components.py — branded "smart components" for Euricom decks.

Each builder draws SIMPLE, native, editable PowerPoint shapes (rounded
rectangles, text boxes, a real table) on a slide returned by
Deck.content_blank(). No groups, no fragile constructions — a consultant
can click any element and edit it afterwards.

All builders share one canvas (below the title) and the brand palette, so
cards, KPIs, timelines and tables look like the same designer made them.
Use fluo-green ONLY as a thin accent (a number, a tick, a short rule) —
never as a fill behind text.

    import components as C
    s = d.content_blank(eyebrow="AANPAK", title="In vier fases")
    C.process_flow(s, ["Discover", "Design", "Build", "Run"])
"""

from __future__ import annotations

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from build_deck import EURICOM, SEMANTIC, CANVAS, _est_lines
import icons

FONT = "Montserrat"
_C = {k: RGBColor.from_string(v) for k, v in EURICOM.items()}  # noqa


def _variant(slide):
    """White-background slides take a subtle shadow; steel-background slides
    should not (per the template's design notes)."""
    return getattr(slide, "_euricom_variant", "white")


def _shadow_for(slide):
    return _variant(slide) != "steel"


# ---------------------------------------------------------------------------
# low-level helpers
# ---------------------------------------------------------------------------
def _txt(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         first=True):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p


def _textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    return tb, tf


def _shadow_subtle(shape):
    """Add a soft, subtle drop shadow (PowerPoint preset-ish, hand-rolled)."""
    spPr = shape._element.spPr
    # remove any inherited shadow first
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(qn("a:outerShdw"),
                           {"blurRad": "90000", "dist": "38100",
                            "dir": "5400000", "rotWithShape": "0"})
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "1D252D"})
    alpha = clr.makeelement(qn("a:alpha"), {"val": "12000"})
    clr.append(alpha)
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def _card(slide, x, y, w, h, fill="white", border="steel", radius=0.06,
          shadow=True, anchor=MSO_ANCHOR.TOP, pad=0.22):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = _C[fill] if fill in _C else RGBColor.from_string(fill)
    sp.line.color.rgb = _C[border] if border in _C else RGBColor.from_string(border)
    sp.line.width = Pt(0.75)
    # small corner radius
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if shadow:
        _shadow_subtle(sp)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(pad))
    return sp


def _fluo_rule(slide, x, y, w=0.42, h=0.045):
    """A short fluo-green accent rule — the recurring subtle accent."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = _C["fluo"]
    sp.line.fill.background()
    if sp.shadow:
        sp.shadow.inherit = False
    return sp


def _tint(hexv, t):
    """Blend a hex colour toward white by fraction t (0..1). t=0 -> original,
    t=1 -> white. Builds a tonal ramp from a single brand colour."""
    r, g, b = (int(hexv[i:i + 2], 16) for i in (0, 2, 4))
    r = round(r + (255 - r) * t)
    g = round(g + (255 - g) * t)
    b = round(b + (255 - b) * t)
    return f"{r:02X}{g:02X}{b:02X}"


def _luma(hexv):
    """Relative luminance to decide readable text colour on a fill."""
    r, g, b = (int(hexv[i:i + 2], 16) for i in (0, 2, 4))
    return 0.299 * r + 0.587 * g + 0.114 * b


def _text_on(hexv):
    """Readable text colour (white or chacoal) for a given background hex."""
    return _C["white"] if _luma(hexv) < 140 else _C["chacoal"]


def _resolve(c, default="midnight"):
    """Resolve a colour given as a brand name (midnight, fluo, ...), a semantic
    name (succes, info, waarschuwing, risico) or a hex string, to an RGBColor.
    Lets KPIs/statements borrow a semantic colour sparingly outside notes."""
    if c is None:
        c = default
    if isinstance(c, RGBColor):
        return c
    key = str(c).lower()
    if key in _C:
        return _C[key]
    if key in SEMANTIC:
        return RGBColor.from_string(SEMANTIC[key][0])
    return RGBColor.from_string(str(c))


# A tonal ramp of midnight tints (deep -> muted teal): on-brand, editable.
# cards(style="tonal") cycles through these in order.
MIDNIGHT_RAMP = [
    EURICOM["midnight"],                  # 014046  deep
    _tint(EURICOM["midnight"], 0.20),
    _tint(EURICOM["midnight"], 0.40),
    _tint(EURICOM["midnight"], 0.60),
]


def _left_bar(slide, x, y, h, color, w=0.07):
    """A thin vertical accent bar along a card's left edge."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color if isinstance(color, RGBColor) \
        else RGBColor.from_string(color)
    sp.line.fill.background()
    if sp.shadow:
        sp.shadow.inherit = False
    return sp


def _fluo_top_cap(slide, x, y, w, h, radius=0.06):
    """A fluo-green rounded base that peeks above a card as a thin top edge
    which follows the rounded corners (the card is drawn on top, shifted down).
    Replaces the old straight left bar that overhung the rounded corner."""
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    base.fill.solid(); base.fill.fore_color.rgb = _C["fluo"]
    base.line.fill.background()
    try:
        base.adjustments[0] = radius
    except Exception:
        pass
    if base.shadow:
        base.shadow.inherit = False
    return base


def _grid_geom(n, columns=None, gap=0.3):
    """Compute card rectangles in the canvas for n items."""
    if columns is None:
        columns = 4 if n >= 4 and n % 4 == 0 else (3 if n != 4 else 2)
        columns = min(columns, n, 4)
    rows = -(-n // columns)
    cw = (CANVAS["w"] - gap * (columns - 1)) / columns
    ch = (CANVAS["h"] - gap * (rows - 1)) / rows
    geoms = []
    for i in range(n):
        r, c = divmod(i, columns)
        x = CANVAS["x"] + c * (cw + gap)
        y = CANVAS["y"] + r * (ch + gap)
        geoms.append((x, y, cw, ch))
    return geoms


# ---------------------------------------------------------------------------
# components
# ---------------------------------------------------------------------------
def cards(slide, items, columns=None, style="cap", accent=None):
    """Cards in a grid. items: list of {"title","body","icon"} or
    (title, body[, icon]) tuples.

    style:
      "cap"    - DEFAULT, the Euricom standard boxed card: white card, subtle
                 border and a rounded fluo-green TOP cap that follows the card's
                 corners (matches the design system's card pages). "fluo" is an
                 alias.
      "plain"  - white cards, steel border, a small fluo accent rule or an icon;
                 the soberest look, no green cap. Use when you want no green.
      "tonal"  - each card a different tint from the MIDNIGHT_RAMP (deep
                 midnight -> muted teal), text colour adapts for contrast.
                 Rich and surprising without extra colour.
      "accent" - one 'hero' midnight card (index given by `accent`, default 0)
                 with a fluo top cap; the rest are standard cap cards.

    The styles are deliberately not combined (no fluo bar on a tonal card):
    one treatment per card keeps the slide calm.
    """
    norm = []
    for it in items:
        if isinstance(it, dict):
            norm.append((it.get("title", ""), it.get("body", ""),
                         it.get("icon")))
        else:
            norm.append((it[0], it[1] if len(it) > 1 else "",
                         it[2] if len(it) > 2 else None))
    n = len(norm)
    columns = columns or min(n, 4) if n != 4 else 2
    columns = min(columns, n)
    rows = -(-n // columns)
    gap = 0.3
    cw = (CANVAS["w"] - gap * (columns - 1)) / columns
    # single row of short cards shouldn't span the whole canvas height
    if rows == 1:
        ch = min(CANVAS["h"], 3.1)
        y0 = CANVAS["y"] + (CANVAS["h"] - ch) / 2
    else:
        ch = (CANVAS["h"] - gap * (rows - 1)) / rows
        y0 = CANVAS["y"]
    hero = (accent if accent is not None else 0) if style == "accent" else -1
    for i, (title, body, icon) in enumerate(norm):
        r, c = divmod(i, columns)
        x = CANVAS["x"] + c * (cw + gap)
        y = y0 + r * (ch + gap)

        # --- per-style fill / border / text colour ------------------------
        if style == "tonal":
            bg = MIDNIGHT_RAMP[i % len(MIDNIGHT_RAMP)]
            fill, border, shadow = bg, bg, False
            title_c = _text_on(bg)
            body_c = title_c
            icon_c = "FFFFFF" if title_c == _C["white"] else EURICOM["midnight"]
            accent_c = "00FF00"
        elif style == "accent" and i == hero:
            bg = EURICOM["midnight"]
            fill, border, shadow = bg, bg, False
            title_c = _C["white"]; body_c = _C["white"]
            icon_c = "FFFFFF"; accent_c = "00FF00"
        elif style in ("fluo", "cap") or (style == "accent" and i != hero):
            # standard Euricom card: subtle border + rounded fluo top cap.
            fill, border, shadow = "white", "steel", _shadow_for(slide)
            title_c = _C["midnight"]; body_c = _C["chacoal"]
            icon_c = EURICOM["midnight"]; accent_c = "00FF00"
        else:  # plain
            fill, border, shadow = "white", "steel", _shadow_for(slide)
            title_c = _C["midnight"]; body_c = _C["chacoal"]
            icon_c = EURICOM["midnight"]; accent_c = "00FF00"

        # fluo / accent-hero cards get a rounded fluo TOP cap (corner-matching):
        # draw the green rounded base, then the card slightly lower so only a
        # thin, rounded green edge shows along the top.
        use_cap = style in ("fluo", "cap", "accent")
        cap_t = 0.055
        radius = 0.06
        if use_cap:
            _fluo_top_cap(slide, x, y, cw, ch, radius=radius)
            card_y, card_h = y + cap_t, ch - cap_t
        else:
            card_y, card_h = y, ch

        card = _card(slide, x, card_y, cw, card_h, fill=fill, border=border,
                     anchor=MSO_ANCHOR.TOP, pad=0.22, shadow=shadow, radius=radius)

        if icon:
            icons.draw(slide, icon, x + 0.22, card_y + 0.30, size=0.52,
                       color=icon_c)
            card.text_frame.margin_top = Inches(1.04)
        elif not use_cap:
            _fluo_rule(slide, x + 0.22, y + 0.22, w=0.32, h=0.04)
            card.text_frame.margin_top = Inches(0.5)
        else:
            card.text_frame.margin_top = Inches(0.5)

        tf = card.text_frame
        if title:
            _txt(tf, title, 14, title_c, bold=True, first=True)
            if body:
                _txt(tf, body, 11, body_c, first=False)
        elif body:
            _txt(tf, body, 11, body_c, first=True)
    return slide


def icon_list(slide, items):
    """A rich vertical list (an alternative to bullets): each row has an icon,
    a bold title and a short description. items: list of {"icon","title",
    "body"} or (icon, title, body) tuples. Fills the whole content area, so do
    not also place a note on the same slide (put the note on its own slide)."""
    norm = []
    for it in items:
        if isinstance(it, dict):
            norm.append((it.get("icon"), it.get("title", ""), it.get("body", "")))
        else:
            norm.append((it[0] if len(it) > 0 else None,
                         it[1] if len(it) > 1 else "",
                         it[2] if len(it) > 2 else ""))
    n = len(norm)
    gap = 0.28
    rh = min((CANVAS["h"] - gap * (n - 1)) / n, 1.5)
    total = rh * n + gap * (n - 1)
    y0 = CANVAS["y"] + max(0, (CANVAS["h"] - total) / 2)
    icon_box = min(rh * 0.8, 0.7)
    text_x = CANVAS["x"] + icon_box + 0.45
    text_w = CANVAS["w"] - icon_box - 0.45
    for i, (icon, title, body) in enumerate(norm):
        y = y0 + i * (rh + gap)
        if icon:
            icons.draw(slide, icon, CANVAS["x"], y + (rh - icon_box) / 2,
                       size=icon_box, color=EURICOM["midnight"])
        else:
            _fluo_rule(slide, CANVAS["x"], y + rh / 2 - 0.02, w=0.42, h=0.045)
        _, tf = _textbox(slide, text_x, y, text_w, rh, anchor=MSO_ANCHOR.MIDDLE)
        if title:
            _txt(tf, title, 15, _C["midnight"], bold=True, first=True)
            if body:
                _txt(tf, body, 11.5, _C["chacoal"], first=False)
        elif body:
            _txt(tf, body, 12, _C["chacoal"], first=True)
        if i < n - 1:                       # subtle divider with a fluo lead
            dy = y + rh + gap / 2
            ln = slide.shapes.add_connector(2, Inches(CANVAS["x"]),
                                            Inches(dy),
                                            Inches(CANVAS["x"] + CANVAS["w"]),
                                            Inches(dy))
            ln.line.color.rgb = _C["steel"]; ln.line.width = Pt(0.5)
            if ln.shadow: ln.shadow.inherit = False
            _fluo_rule(slide, CANVAS["x"], dy - 0.02, w=0.5, h=0.045)
    return slide


def kpi_row(slide, items):
    """Big-number stat callouts. items: list of {"value","label","icon","color"}
    or (value, label) tuples. Number in Midnight by default; pass a semantic
    name ("succes"/"risico"/...) or brand name/hex as "color" to tint a figure
    (use sparingly). Small muted label below, optional icon above. Centred."""
    norm = []
    for it in items:
        if isinstance(it, dict):
            norm.append((it.get("value", ""), it.get("label", ""),
                         it.get("icon"), it.get("color")))
        else:
            norm.append((it[0], it[1] if len(it) > 1 else "", None, None))
    n = len(norm)
    gap = 0.3
    cw = (CANVAS["w"] - gap * (n - 1)) / n
    block_h = 2.1
    y = CANVAS["y"] + max(0, (CANVAS["h"] - block_h) / 2)
    for i, (value, label, icon, color) in enumerate(norm):
        x = CANVAS["x"] + i * (cw + gap)
        num_c = _resolve(color)
        if icon:
            icons.draw(slide, icon, x, y, size=0.6,
                       color="%02X%02X%02X" % (num_c[0], num_c[1], num_c[2]))
            ny = y + 0.85
        else:
            ny = y + 0.3
        _, tf = _textbox(slide, x, ny, cw, 1.0)
        _txt(tf, str(value), 50, num_c, bold=True)
        # subtle fluo accent under every figure (also when an icon is present)
        _fluo_rule(slide, x, ny + 0.95, w=0.46, h=0.05)
        _, tf2 = _textbox(slide, x, ny + 1.12, cw, 0.8)
        _txt(tf2, label, 12.5, _C["light_mid"])
    return slide


def statement(slide, text, sub=None, top=None, color=None):
    """One strong takeaway per slide: a large bold statement with a fluo accent
    rule, vertically centred in the content area. Optional smaller sub-line.
    Pass `top` (inches) to anchor it there instead of centring — useful when a
    `note` or other element shares the slide below it."""
    n = len(text or "")
    size = 40 if n <= 60 else 34 if n <= 110 else 28 if n <= 170 else 24
    lines = _est_lines(text, CANVAS["w"] * 0.88, size, bold=True)
    line_h = size / 72 * 1.2
    text_h = lines * line_h
    sub_h = 0.5 if sub else 0.0
    block_h = 0.06 + 0.28 + text_h + (0.2 + sub_h if sub else 0)
    if top is not None:
        y = top
    else:
        y = CANVAS["y"] + max(0, (CANVAS["h"] - block_h) / 2)
    _fluo_rule(slide, CANVAS["x"], y, w=0.6, h=0.06)
    _, tf = _textbox(slide, CANVAS["x"], y + 0.28, CANVAS["w"] * 0.88,
                     text_h + 0.2)
    _txt(tf, text, size, _resolve(color), bold=True)
    if sub:
        _, tf2 = _textbox(slide, CANVAS["x"], y + 0.28 + text_h + 0.2,
                          CANVAS["w"] * 0.8, sub_h)
        _txt(tf2, sub, 14, _C["light_mid"])
    return slide


def process_flow(slide, steps):
    """Horizontal numbered steps. steps: list of str, or list of
    {"title","body"}."""
    norm = []
    for s in steps:
        if isinstance(s, dict):
            norm.append((s.get("title", ""), s.get("body", "")))
        else:
            norm.append((str(s), ""))
    n = len(norm)
    gap = 0.3
    cw = (CANVAS["w"] - gap * (n - 1)) / n
    y = CANVAS["y"] + 0.7
    h = 1.9
    for i, (title, body) in enumerate(norm):
        x = CANVAS["x"] + i * (cw + gap)
        card = _card(slide, x, y, cw, h)
        tf = card.text_frame
        _txt(tf, f"{i + 1:02d}", 13, _C["midnight"], bold=True, first=True)
        _txt(tf, title, 13, _C["midnight"], bold=True, first=False)
        if body:
            _txt(tf, body, 10.5, _C["chacoal"], first=False)
        # fluo connector between steps
        if i < n - 1:
            _fluo_rule(slide, x + cw + 0.02, y + h / 2 - 0.02,
                       w=gap - 0.04, h=0.045)
    return slide


def timeline(slide, items):
    """Horizontal timeline. items: list of {"date","title","body"} or
    (date, title, body) tuples."""
    norm = []
    for it in items:
        if isinstance(it, dict):
            norm.append((it.get("date", ""), it.get("title", ""),
                         it.get("body", "")))
        else:
            norm.append((it[0], it[1] if len(it) > 1 else "",
                         it[2] if len(it) > 2 else ""))
    n = len(norm)
    line_y = CANVAS["y"] + 0.55
    # base line in steel
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(CANVAS["x"]), Inches(line_y),
                                  Inches(CANVAS["w"]), Inches(0.03))
    base.fill.solid(); base.fill.fore_color.rgb = _C["steel"]
    base.line.fill.background()
    if base.shadow: base.shadow.inherit = False
    step = CANVAS["w"] / n
    for i, (date, title, body) in enumerate(norm):
        cx = CANVAS["x"] + step * i + step / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(cx - 0.09), Inches(line_y - 0.075),
                                     Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = _C["fluo"]
        dot.line.color.rgb = _C["midnight"]; dot.line.width = Pt(1)
        if dot.shadow: dot.shadow.inherit = False
        _, tf = _textbox(slide, cx - step / 2 + 0.1, line_y + 0.3,
                         step - 0.2, 2.2)
        if date:  _txt(tf, str(date), 11, _C["light_mid"], bold=True)
        if title: _txt(tf, title, 13, _C["midnight"], bold=True,
                       first=not bool(date))
        if body:  _txt(tf, body, 10.5, _C["chacoal"], first=False)
    return slide


def comparison(slide, left, right, left_title=None, right_title=None):
    """Two-column comparison. left/right: list of str. Optional headers."""
    gap = 0.4
    cw = (CANVAS["w"] - gap) / 2
    for col, (title, rows) in enumerate(
            [(left_title, left), (right_title, right)]):
        x = CANVAS["x"] + col * (cw + gap)
        card = _card(slide, x, CANVAS["y"], cw, CANVAS["h"],
                     shadow=_shadow_for(slide))
        _fluo_rule(slide, x + 0.22, CANVAS["y"] + 0.26, w=0.42, h=0.05)
        card.text_frame.margin_top = Inches(0.62)
        tf = card.text_frame
        if title:
            tp = _txt(tf, title, 15, _C["midnight"], bold=True, first=True)
            tp.space_after = Pt(10)
            start_first = False
        else:
            start_first = True
        for j, row in enumerate(rows):
            p = _txt(tf, str(row), 12, _C["chacoal"],
                     first=(start_first and j == 0))
            p.space_after = Pt(6)
    return slide


def maturity(slide, levels):
    """Ascending maturity bars. levels: ordered list of str (low->high) or
    {"title","body"}."""
    norm = [(l["title"], l.get("body", "")) if isinstance(l, dict)
            else (str(l), "") for l in levels]
    n = len(norm)
    gap = 0.25
    cw = (CANVAS["w"] - gap * (n - 1)) / n
    max_h = CANVAS["h"] - 0.6
    min_h = max_h * 0.42
    for i, (title, body) in enumerate(norm):
        h = min_h + (max_h - min_h) * (i / max(1, n - 1))
        x = CANVAS["x"] + i * (cw + gap)
        y = CANVAS["y"] + (CANVAS["h"] - h)
        fill = "midnight" if i == n - 1 else ("light_mid" if i >= n - 2 else "steel")
        bar = _card(slide, x, y, cw, h, fill=fill, border=fill,
                    radius=0.04, shadow=False)
        tf = bar.text_frame
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        txt_color = _C["white"] if i >= n - 2 else _C["midnight"]
        _txt(tf, title, 13, txt_color, bold=True, first=True)
        if body:
            _txt(tf, body, 10, txt_color, first=False)
    return slide


def data_table(slide, headers, rows, x=None, y=None, w=None, h=None,
               total_row=False):
    """Branded data table: Midnight header (white bold), zebra body, steel
    borders. A real, editable PowerPoint table."""
    x = CANVAS["x"] if x is None else x
    y = CANVAS["y"] if y is None else y
    w = CANVAS["w"] if w is None else w
    nrows = len(rows) + 1
    ncols = len(headers)
    h = (min(0.55 * nrows + 0.1, CANVAS["h"]) if h is None else h)
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                                  Inches(w), Inches(h))
    tbl = gtbl.table
    tbl.first_row = True
    tbl.horz_banding = True

    def cell_text(cell, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold
        f.color.rgb = color

    def shade(cell, hexv):
        tcPr = cell._tc.get_or_add_tcPr()
        for el in tcPr.findall(qn("a:fill")) + tcPr.findall(qn("a:solidFill")):
            tcPr.remove(el)
        fill = tcPr.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": hexv})
        fill.append(clr); tcPr.append(fill)

    for c, head in enumerate(headers):
        cell = tbl.cell(0, c)
        shade(cell, EURICOM["midnight"])
        cell_text(cell, head, 10.5, _C["white"], bold=True)
    for ri, row in enumerate(rows, start=1):
        is_total = total_row and ri == len(rows)
        bg = "CBD9DA" if is_total else (EURICOM["light_steel"] if ri % 2 == 0
                                        else "FFFFFF")
        for c in range(ncols):
            cell = tbl.cell(ri, c)
            shade(cell, bg)
            val = row[c] if c < len(row) else ""
            cell_text(cell, val, 10, _C["chacoal"], bold=is_total)
    return gtbl


def _content_bottom(slide):
    """Lowest y (inches) occupied by real content sitting in the canvas area,
    ignoring placeholders above it (title/eyebrow) and the page number below."""
    bottom = CANVAS["y"]
    for sh in slide.shapes:
        try:
            t = sh.top / 914400
            h = (sh.height or 0) / 914400
        except Exception:
            continue
        if t is None:
            continue
        # only shapes that actually sit within the canvas band
        if t < CANVAS["y"] + CANVAS["h"] - 0.05 and (t + h) > CANVAS["y"] + 0.05:
            bottom = max(bottom, t + h)
    return bottom


def _next_note_y(slide):
    """Auto-placement for notes: stack below the previous note, or below any
    content already in the canvas, so notes never land on top of each other
    or over a builder. Falls back to the canvas top on an empty slide."""
    cur = getattr(slide, "_euricom_note_y", None)
    if cur is not None:
        return cur
    b = _content_bottom(slide)
    return b + 0.2 if b > CANVAS["y"] + 0.05 else CANVAS["y"]


# Default icon + fallback label per semantic kind (mirrors the Word notes).
_NOTE_ICON = {"succes": "check", "waarschuwing": "flag",
              "risico": "shield", "info": "bulb"}
_NOTE_LABEL = {"succes": "Succes", "waarschuwing": "Waarschuwing",
               "risico": "Risico", "info": "Info"}


def note(slide, kind, title=None, body=None, x=None, y=None, w=None,
         icon=True):
    """A light, slide-native semantic call-out (succes / waarschuwing / risico /
    info): a slim rounded semantic bar + icon + bold label + text on the slide
    background (no coloured fill box — that reads too "documenty" on a slide).
    `title` defaults to the type label; pass `icon=False` to drop the icon, or
    `icon="name"` to override. Use sparingly: only for a genuine decision or
    risk. For an informational aside, prefer a `statement` with one
    semantic-coloured word instead of a note.

    Placement is automatic: when `y` is omitted the note stacks below the
    previous note (and below any content already on the slide), so notes never
    overlap each other or a builder. If there is no room left in the canvas the
    note is SKIPPED (returns None) rather than pushed off-slide — notes are
    meant to be sparing. Pass an explicit `y=` to force a position."""
    kind = kind.lower()
    fg, bg = SEMANTIC.get(kind, SEMANTIC["info"])
    title = title or _NOTE_LABEL.get(kind, "Info")
    x = CANVAS["x"] if x is None else x
    auto_y = y is None
    y = _next_note_y(slide) if auto_y else y
    w = CANVAS["w"] if w is None else w
    fgc = RGBColor.from_string(fg)

    ico = _NOTE_ICON.get(kind) if icon is True else (icon or None)
    bar_w, ico_box = 0.09, 0.42
    icon_x = x + 0.30
    text_l = icon_x + (ico_box + 0.22 if ico else 0)
    text_w = (x + w) - text_l

    # auto-height from title + body line counts
    lines = _est_lines(title, text_w, 12, bold=True)
    if body:
        lines += _est_lines(body, text_w, 11)
    h = max(0.95, 0.34 + lines * 0.26)
    if auto_y and y + h > CANVAS["y"] + CANVAS["h"] + 0.05:
        return None          # no room left in the canvas -> skip (notes are
                             # sparing; better none than overflowing the slide)
    slide._euricom_note_y = y + h + 0.2     # next note stacks below this one

    # slim rounded semantic bar (pill ends), no fill box: light and slide-native
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(bar_w), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = fgc
    bar.line.fill.background()
    try: bar.adjustments[0] = 0.5       # fully rounded (pill) ends
    except Exception: pass
    if bar.shadow: bar.shadow.inherit = False

    if ico:
        icons.draw(slide, ico, icon_x, y + (h - ico_box) / 2,
                   size=ico_box, color=fg)

    _, tf = _textbox(slide, text_l, y, text_w, h, anchor=MSO_ANCHOR.MIDDLE)
    _txt(tf, title, 12, fgc, bold=True, first=True)
    if body:
        _txt(tf, body, 11, _C["chacoal"], first=False)
    return bar


# ---------------------------------------------------------------------------
# images
# ---------------------------------------------------------------------------
def _place_cover(slide, path, x, y, w, h):
    """Add a picture sized to (w,h) and cover-cropped (no distortion): the
    image fills the box, excess is trimmed equally on the long axis."""
    from PIL import Image as _PILImage
    iw, ih = _PILImage.open(path).size
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    box_ar = w / h
    img_ar = iw / ih
    if img_ar > box_ar:                      # image too wide -> crop sides
        crop = (1 - box_ar / img_ar) / 2
        pic.crop_left = crop; pic.crop_right = crop
    elif img_ar < box_ar:                    # image too tall -> crop top/bottom
        crop = (1 - img_ar / box_ar) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic


def _pic_geom(pic, prst, adj=None):
    """Reshape a picture to a preset geometry (e.g. 'roundRect', 'ellipse')."""
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    geom = spPr.makeelement(qn("a:prstGeom"), {"prst": prst})
    av = geom.makeelement(qn("a:avLst"), {})
    if adj is not None:
        gd = av.makeelement(qn("a:gd"),
                            {"name": "adj", "fmla": "val %d" % int(adj * 100000)})
        av.append(gd)
    geom.append(av)
    xfrm = spPr.find(qn("a:xfrm"))            # geometry must follow xfrm
    if xfrm is not None:
        xfrm.addnext(geom)
    else:
        spPr.insert(0, geom)
    return pic


def image_grid(slide, paths, cols=2, rows=None, gap=0.2, radius=0.05,
               x=None, y=None, w=None, h=None, shadow=True):
    """A grid of rounded, cover-cropped photos filling a rect (default CANVAS).
    Pass x/y/w/h to place it in part of the slide (e.g. the right half)."""
    x = CANVAS["x"] if x is None else x
    y = CANVAS["y"] if y is None else y
    w = CANVAS["w"] if w is None else w
    h = CANVAS["h"] if h is None else h
    n = len(paths)
    rows = rows or -(-n // cols)
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rows - 1)) / rows
    for i, p in enumerate(paths):
        r, c = divmod(i, cols)
        px = x + c * (cw + gap)
        py = y + r * (ch + gap)
        pic = _place_cover(slide, p, px, py, cw, ch)
        _pic_geom(pic, "roundRect", adj=radius)
        if shadow and _shadow_for(slide):
            _shadow_subtle(pic)
    return slide


def team(slide, members, columns=None):
    """A row of people: circular photo, bold name, muted role, short body.
    members: list of {"photo","name","role","body"}. Photos are cover-cropped
    to a circle. Fills the content area as N equal columns."""
    n = len(members)
    columns = columns or n
    gap = 0.4
    cw = (CANVAS["w"] - gap * (columns - 1)) / columns
    photo = min(cw * 0.55, 1.7)
    y = CANVAS["y"] + 0.05
    for i, m in enumerate(members):
        x = CANVAS["x"] + i * (cw + gap)
        if m.get("photo"):
            pic = _place_cover(slide, m["photo"], x, y, photo, photo)
            _pic_geom(pic, "ellipse")
            if _shadow_for(slide):
                _shadow_subtle(pic)
        ty = y + photo + 0.28
        _, tf = _textbox(slide, x, ty, cw, 0.4)
        _txt(tf, m.get("name", ""), 15, _C["midnight"], bold=True)
        _, tf2 = _textbox(slide, x, ty + 0.42, cw, 0.35)
        _txt(tf2, m.get("role", ""), 12.5, _C["light_mid"])
        if m.get("body"):
            _, tf3 = _textbox(slide, x, ty + 0.92, cw, 1.6)
            _txt(tf3, m["body"], 11.5, _C["chacoal"])
    return slide
