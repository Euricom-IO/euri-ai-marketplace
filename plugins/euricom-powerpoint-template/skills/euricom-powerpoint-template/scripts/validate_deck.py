#!/usr/bin/env python3
"""
validate_deck.py - structural QA for a generated Euricom deck. Run BEFORE
presenting. Checks the things that silently break a deck: a slide that lost
its branded layout, a leftover content-area rectangle or placeholder prompt
text, a stripped theme, leftover sections, or an empty deck. Visual QA
(overflow, overlap) still needs a render pass - see SKILL.md.

    python validate_deck.py /mnt/user-data/outputs/Deck-v01.pptx
"""
import sys
import zipfile
from pptx import Presentation
from pptx.oxml.ns import qn

EURICOM_LAYOUTS = {"10_Cover_Dark", "20_Content_Steal", "30_Content_White",
                   "30_Section_Dark", "40_Quote"}
PROMPT_MARKERS = ["click to edit", "klik om", "lorem", "ipsum",
                  "tijdelijke aanduiding", "your text"]
_P14_SECTIONLST = "{http://schemas.microsoft.com/office/powerpoint/2010/main}sectionLst"


def _line_hex(shape):
    try:
        if shape.line.color is not None and shape.line.color.type is not None:
            return str(shape.line.color.rgb)
    except Exception:
        pass
    return None


def validate(path):
    errors, warnings, info = [], [], []
    try:
        prs = Presentation(path)
    except Exception as e:
        return [f"File will not open: {e}"], [], []

    n = len(prs.slides._sldIdLst)
    info.append(f"{n} slide(s)")
    if n == 0:
        errors.append("Deck has no slides.")

    used = {}
    no_notes = []
    for i, s in enumerate(prs.slides, 1):
        name = s.slide_layout.name
        used[name] = used.get(name, 0) + 1
        if name not in EURICOM_LAYOUTS:
            warnings.append(f"Slide {i} uses a non-Euricom layout '{name}'.")
        # speaker notes are optional (recommended, not required)
        has_notes = (s.has_notes_slide
                     and (s.notes_slide.notes_text_frame.text or "").strip())
        if not has_notes:
            no_notes.append(i)
        for sh in s.shapes:
            if not sh.is_placeholder and _line_hex(sh) == "FF0000":
                errors.append(f"Slide {i}: leftover red content-area rectangle "
                              f"- it must be removed from output.")
        # leftover example/prompt text in title/eyebrow
        for sh in s.shapes:
            if sh.has_text_frame:
                low = sh.text_frame.text.strip().lower()
                if low and any(mk in low for mk in PROMPT_MARKERS):
                    warnings.append(f"Slide {i}: possible leftover prompt text "
                                    f"{sh.text_frame.text.strip()[:40]!r}.")
    info.append("layout usage: " +
                ", ".join(f"{k}x{v}" for k, v in used.items()))
    if no_notes:
        info.append(f"{len(no_notes)} slide(s) without speaker notes "
                    f"(optional - notes are recommended, not required).")

    if prs.part._element.find(".//" + _P14_SECTIONLST) is not None:
        warnings.append("Section list still present - base/sample/library "
                        "sections should be stripped from output.")

    try:
        z = zipfile.ZipFile(path)
        th = z.read("ppt/theme/theme1.xml").decode("utf-8", "ignore")
        if "Montserrat" not in th:
            errors.append("Theme font Montserrat missing - template theme lost.")
        if "014046" not in th:
            errors.append("Brand colour Midnight (014046) missing from theme.")
        names = z.namelist()
        dupes = sorted({x for x in names if names.count(x) > 1})
        if dupes:
            errors.append(f"Duplicate package parts: {dupes[:5]}")
    except Exception as e:
        warnings.append(f"Could not inspect package: {e}")

    return errors, warnings, info


def main():
    if len(sys.argv) < 2:
        print("usage: python validate_deck.py <deck.pptx>")
        sys.exit(2)
    errors, warnings, info = validate(sys.argv[1])
    for x in info:
        print(f"  - {x}")
    for w in warnings:
        print(f"  ! {w}")
    for e in errors:
        print(f"  X {e}")
    if errors:
        print("RESULT: FAIL")
        sys.exit(1)
    print("RESULT: PASS" + (" (with warnings)" if warnings else ""))


if __name__ == "__main__":
    main()
