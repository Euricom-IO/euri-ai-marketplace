#!/usr/bin/env python3
"""
build_deck.py - Euricom PowerPoint design-system engine (v05 template).

The template is the single source of truth and is self-documenting. It holds
three native PowerPoint sections:

  * "Base Library"       - one filled, annotated example per master layout
                           (cover / content-steal / content-white / section /
                           quote). Each slide's speaker notes explain how to
                           fill it.
  * "Components Library" - richer, composed reusable slides (e.g. the agenda
                           table). Notes describe each component's purpose.
  * "Sample Presentation"- scratch space, ignored.

ONE consistent model: clone a reference slide and refill it. Standard slides
clone the matching Base Library example; richer slides clone a Components
Library entry. Cloning inherits everything the designer baked in - exact
placeholder formatting, the automatic page number, table styling - so output
always matches a hand-made deck and stays fully editable (real placeholders,
native bullet lists, simple native shapes).

A red rectangle on the content example marks the content area for free-drawn
components; the engine reads its geometry (CANVAS) and removes it from output.

Workflow: ALWAYS read the catalogue first (library()/base_examples()); reuse a
component when one fits; otherwise build on a content example and draw within
CANVAS. On save, the Base/Components/Sample originals and the section list are
stripped, so the deliverable contains only generated slides.
"""

from __future__ import annotations

import copy
import glob
import math
import os
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# --- Brand palette (single source of truth; mirrors the Word template) ------
EURICOM = {
    "chacoal": "1D252D", "white": "FFFFFF", "midnight": "014046",
    "light_steel": "F1F5F6", "fluo": "00FF00", "light_mid": "809FA2",
    "steel": "CBD9DA",
}
SEMANTIC = {
    # succes uses Midnight (not the old teal-green 30CBB1) on a light-steel
    # ground, mirroring the "positive/tip" note in the Euricom Word template.
    "succes": ("014046", "F1F5F6"), "waarschuwing": ("E9AB0C", "FEF9EC"),
    "risico": ("E80F0F", "FEF0F0"), "info": ("5C7B7E", "F1F5F6"),
}

_LAYOUT_NAMES = {
    "cover": "10_Cover_Dark", "steel": "20_Content_Steal",
    "white": "30_Content_White", "section": "30_Section_Dark",
    "quote": "40_Quote",
}
_LAYOUT_FALLBACK_IDX = {"cover": 0, "steel": 1, "white": 2, "section": 3, "quote": 4}

PH = {
    "cover": {"eyebrow": 11, "title": 12, "subtitle": 13, "intro": 14,
              "date": 15, "author": 16},
    "section": {"eyebrow": 11, "title": 12, "description": 14},
    "content": {"eyebrow": 11, "title": 0, "body": 12},
    "quote": {"quote": 12, "author": 16, "role": 17},
}

# Default content area; overridden at runtime by the red rectangle if present.
CANVAS = {"x": 0.76, "y": 1.94, "w": 11.81, "h": 5.08}

_BASE_SECTION = "base library"
_LIBRARY_SECTION = "components library"
_SAMPLE_SECTION = "sample presentation"
_CONTENT_AREA_COLOR = "FF0000"

_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_P14_SECTION = f"{{{_P14}}}section"
_P14_SLDID = f"{{{_P14}}}sldId"
_P14_SECTIONLST = f"{{{_P14}}}sectionLst"


# --- text helpers -----------------------------------------------------------
def _est_lines(text, box_width_in, font_pt, bold=False):
    if not text:
        return 1
    factor = 0.62 if bold else 0.55
    cpl = max(8, int(box_width_in * 72 / (font_pt * factor)))
    return max(1, math.ceil(len(text) / cpl * 1.08))


def _autofit_shrink(text_frame):
    from pptx.enum.text import MSO_AUTO_SIZE
    try:
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _no_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _find_template(explicit=None):
    if explicit and os.path.exists(explicit):
        return explicit
    ups = [u for u in glob.glob("/mnt/user-data/uploads/*.pptx")
           if "euricom" in os.path.basename(u).lower()
           or "template" in os.path.basename(u).lower()]
    if ups:
        return max(ups, key=os.path.getmtime)   # newest version wins
    shipped = Path(__file__).resolve().parent.parent / "assets" / \
        "Euricom_PowerPoint_Template.pptx"
    if shipped.exists():
        return str(shipped)
    hits = glob.glob("/mnt/skills/**/Euricom_PowerPoint_Template*.pptx",
                     recursive=True)
    if hits:
        return hits[0]
    raise FileNotFoundError("Euricom PowerPoint template not found.")


def _line_hex(shape):
    """Read a shape's line colour from the XML, WITHOUT going through
    python-pptx' `shape.line.color` — that accessor mutates a `<a:ln><a:noFill/>`
    into an (empty) `<a:ln><a:solidFill/>`, which then renders as a stray black
    border on every cloned shape it touches. Returns the hex of an explicit
    line solidFill, else None."""
    try:
        spPr = shape._element.spPr
    except Exception:
        return None
    if spPr is None:
        return None
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return None
    sf = ln.find(qn("a:solidFill"))
    if sf is None:
        return None
    clr = sf.find(qn("a:srgbClr"))
    return clr.get("val") if clr is not None else None


_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _duplicate_slide(prs, src_slide):
    """Duplicate a slide within the same presentation, keeping its layout and
    copying its content. Image/other relationships are re-added via the public
    `relate_to` API (python-pptx version-proof) and the relationship ids in the
    copied shape XML are remapped to the new ids so pictures still resolve."""
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    new_spTree = new_slide.shapes._spTree
    copied = []
    for child in list(src_slide.shapes._spTree):
        if child.tag.split('}')[-1] in ('sp', 'pic', 'graphicFrame', 'grpSp',
                                         'cxnSp'):
            el = copy.deepcopy(child)
            new_spTree.append(el)
            copied.append(el)
    # re-add relationships (images etc.) and map old rId -> new rId
    id_map = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype.split('/')[-1] in ('slideLayout', 'notesSlide'):
            continue
        try:
            if rel.is_external:
                new_rId = new_slide.part.relate_to(rel.target_ref, rel.reltype,
                                                   is_external=True)
            else:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            continue
        if new_rId and new_rId != rId:
            id_map[rId] = new_rId
    if id_map:
        for el in copied:
            for node in el.iter():
                for k, v in list(node.attrib.items()):
                    if k.startswith("{" + _REL_NS + "}") and v in id_map:
                        node.set(k, id_map[v])
    return new_slide


class Deck:
    def __init__(self, template=None):
        self.template_path = _find_template(template)
        self.prs = Presentation(self.template_path)
        self._layouts = {k: self._layout(n, _LAYOUT_FALLBACK_IDX[k])
                         for k, n in _LAYOUT_NAMES.items()}
        self._original_rids = {sid.get(qn("r:id"))
                               for sid in self.prs.slides._sldIdLst}
        self._sections = self._read_sections()
        self._base = self._index_base()          # role -> example slide
        self._catalogue = self._index_library()   # rich components
        self.canvas = self._read_content_area()    # from the red rectangle
        CANVAS.update(self.canvas)   # share with components.py (same dict)

    # -- layout / placeholder internals -----------------------------------
    def _layout(self, name, fallback_idx):
        for lay in self.prs.slide_layouts:
            if lay.name == name:
                return lay
        return self.prs.slide_layouts[fallback_idx]

    @staticmethod
    def _ph(slide, idx):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _layout_geom(self, key, idx):
        for ph in self._layouts[key].placeholders:
            if ph.placeholder_format.idx == idx:
                return ph.left, ph.top, ph.width, ph.height
        return None

    @staticmethod
    def _place(ph, left=None, top=None, width=None, height=None):
        if left is not None:   ph.left = left
        if top is not None:    ph.top = top
        if width is not None:  ph.width = width
        if height is not None: ph.height = height

    # -- sections / catalogue ---------------------------------------------
    def _read_sections(self):
        """Map slide-id -> section name (lower-cased)."""
        out = {}
        for sec in self.prs.part._element.iter(_P14_SECTION):
            name = (sec.get("name") or "").strip().lower()
            for sid in sec.iter(_P14_SLDID):
                out[sid.get("id")] = name
        return out

    def _section_of(self, position):
        sid = self.prs.slides._sldIdLst[position].get("id")
        return self._sections.get(sid)

    def _notes(self, slide):
        if not slide.has_notes_slide:
            return ""
        tf = slide.notes_slide.notes_text_frame
        return (tf.text or "").strip() if tf is not None else ""

    def _index_base(self):
        base = {}
        name_to_role = {v: k for k, v in _LAYOUT_NAMES.items()}
        for pos, slide in enumerate(self.prs.slides):
            if self._section_of(pos) != _BASE_SECTION:
                continue
            role = name_to_role.get(slide.slide_layout.name)
            if role and role not in base:
                base[role] = slide
        return base

    def _index_library(self):
        cat = []
        for pos, slide in enumerate(self.prs.slides):
            if self._section_of(pos) != _LIBRARY_SECTION:
                continue
            notes = self._notes(slide)
            title = ""
            t = self._ph(slide, 0)
            if t is not None and t.has_text_frame:
                title = t.text_frame.text.strip()
            label = (notes.splitlines()[0].strip() if notes else title) \
                or f"component {pos}"
            cat.append({"index": pos, "_slide": slide, "title": title,
                        "notes": notes, "label": label,
                        "layout": slide.slide_layout.name})
        return cat

    def _read_content_area(self):
        for slide in self._base.values():
            for sh in slide.shapes:
                if not sh.is_placeholder and _line_hex(sh) == _CONTENT_AREA_COLOR:
                    return {"x": round(sh.left / 914400, 3),
                            "y": round(sh.top / 914400, 3),
                            "w": round(sh.width / 914400, 3),
                            "h": round(sh.height / 914400, 3)}
        return dict(CANVAS)

    def base_examples(self):
        """Roles that have a Base Library example, with their notes (the
        per-layout filling instructions)."""
        return {role: self._notes(s) for role, s in self._base.items()}

    def library(self):
        return [{k: c[k] for k in ("index", "title", "label", "notes", "layout")}
                for c in self._catalogue]

    # -- refill helpers ----------------------------------------------------
    @staticmethod
    def _strip_content_area(slide):
        for sh in list(slide.shapes):
            if not sh.is_placeholder and _line_hex(sh) == _CONTENT_AREA_COLOR:
                sh._element.getparent().remove(sh._element)

    def _refill(self, slide, idx, text):
        """Overwrite a placeholder's text while preserving its (inherited or
        run-level) formatting: edit the first run in place, clear the rest.
        If the cloned placeholder box was collapsed in the example (an empty
        field can shrink to ~0 width/height), drop the override so it inherits
        the layout's real geometry."""
        ph = self._ph(slide, idx)
        if ph is None or text is None:
            return
        self._uncollapse(ph)
        tf = ph.text_frame
        paras = tf.paragraphs
        p0 = paras[0]
        for extra in paras[1:]:
            extra._p.getparent().remove(extra._p)
        if p0.runs:
            p0.runs[0].text = str(text)
            for r in p0.runs[1:]:
                r._r.getparent().remove(r._r)
        else:
            p0.add_run().text = str(text)

    @staticmethod
    def _uncollapse(ph):
        """Remove a near-zero-size xfrm override so the placeholder inherits
        its proper position/size from the layout."""
        spPr = ph._element.find(qn("p:spPr"))
        if spPr is None:
            return
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is None:
            return
        ext = xfrm.find(qn("a:ext"))
        if ext is None:
            return
        cx = int(ext.get("cx", "0"))
        cy = int(ext.get("cy", "0"))
        if cx < 457200 or cy < 45720:     # < 0.5" wide or < 0.05" tall
            spPr.remove(xfrm)

    def _drop(self, slide, idx):
        ph = self._ph(slide, idx)
        if ph is not None:
            ph._element.getparent().remove(ph._element)

    def _clear(self, slide, idx):
        """Empty a placeholder's text (keep the placeholder)."""
        ph = self._ph(slide, idx)
        if ph is None:
            return
        for p in ph.text_frame.paragraphs[1:]:
            p._p.getparent().remove(p._p)
        p0 = ph.text_frame.paragraphs[0]
        for r in list(p0.runs):
            r._r.getparent().remove(r._r)

    # -- factory: clone a base example (preferred) or instantiate layout ---
    def _new(self, role):
        base = self._base.get(role)
        if base is not None:
            s = _duplicate_slide(self.prs, base)
            self._strip_content_area(s)
            return s
        # fallback: bare layout + cloned slide-number placeholder
        key = role
        s = self.prs.slides.add_slide(self._layouts[key])
        layout = self._layouts[key]
        for ph in layout.placeholders:
            node = ph._element.find(".//" + qn("p:ph"))
            if node is not None and node.get("type") == "sldNum":
                s.shapes._spTree.append(copy.deepcopy(ph._element))
                break
        return s

    # -- standard slides ---------------------------------------------------
    def cover(self, title, eyebrow=None, subtitle=None, intro=None,
              date=None, author=None, notes=None):
        s = self._new("cover")
        m = PH["cover"]
        self._refill(s, m["eyebrow"], eyebrow.upper()) if eyebrow else self._drop(s, m["eyebrow"])
        self._refill(s, m["title"], title)
        for fld, val in (("subtitle", subtitle), ("intro", intro),
                         ("date", date), ("author", author)):
            if val:
                self._refill(s, m[fld], val)
            else:
                self._drop(s, m[fld])
        self._layout_cover(s, title, subtitle, intro)
        self.set_notes(s, notes)
        return s

    def _layout_cover(self, s, title, subtitle, intro):
        """Stack eyebrow -> title -> subtitle -> intro from a fixed top, sizing
        each box to its real line count and pushing the date/author block down,
        so a long title or tagline never overlaps the text below it. Long
        titles/taglines also step down in size."""
        m = PH["cover"]
        EMU = 914400

        def _set_size(idx, pt):
            ph = self._ph(s, idx)
            if ph is None:
                return
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(pt)

        def _geom(idx):
            return self._layout_geom("cover", idx)

        top = 2.48                       # title's anchor top (from the layout)

        # --- title: 36pt, step down when long; allow up to ~3 lines ---------
        tl = len(title or "")
        t_size = 36 if tl <= 24 else 30 if tl <= 44 else 26
        tg = _geom(m["title"])
        t_w_in = (tg[2] / EMU) if tg else 9.27
        t_lines = min(_est_lines(title, t_w_in, t_size, bold=True), 3)
        t_lh = t_size / 72 * 1.18
        _set_size(m["title"], t_size)
        tph = self._ph(s, m["title"])
        if tph is not None and tg:
            self._place(tph, left=tg[0], top=Inches(top), width=tg[2],
                        height=Inches(max(t_lines * t_lh, tg[3] / EMU)))
        y = top + max(t_lines * t_lh, (tg[3] / EMU) if tg else t_lh)

        # --- subtitle (fluo tagline): 36pt, step down when long -------------
        sub_ph = self._ph(s, m["subtitle"])
        if sub_ph is not None and subtitle:
            sl = len(subtitle)
            s_size = 36 if sl <= 30 else 30 if sl <= 52 else 24
            sg = _geom(m["subtitle"])
            s_w_in = (sg[2] / EMU) if sg else 8.89
            s_lines = min(_est_lines(subtitle, s_w_in, s_size), 3)
            s_lh = s_size / 72 * 1.18
            _set_size(m["subtitle"], s_size)
            y += 0.12
            self._place(sub_ph, left=sg[0], top=Inches(y), width=sg[2],
                        height=Inches(s_lines * s_lh))
            y += s_lines * s_lh

        # --- intro: keep its small steel size; size box to its lines --------
        intro_ph = self._ph(s, m["intro"])
        if intro_ph is not None and intro:
            ig = _geom(m["intro"])
            i_w_in = (ig[2] / EMU) if ig else 5.89
            i_lines = _est_lines(intro, i_w_in, 14)
            i_lh = 14 / 72 * 1.30
            y += 0.18
            self._place(intro_ph, left=ig[0], top=Inches(y), width=ig[2],
                        height=Inches(max(i_lines * i_lh, 0.4)))
            y += max(i_lines * i_lh, 0.4)

        # --- push date/author down if the block grew into them --------------
        dg = _geom(m["date"])
        ag = _geom(m["author"])
        date_top = (dg[1] / EMU) if dg else 5.33
        author_top = (ag[1] / EMU) if ag else 5.55
        shift = max(0.0, (y + 0.3) - date_top)
        if shift > 0:
            dph = self._ph(s, m["date"])
            aph = self._ph(s, m["author"])
            if dph is not None and dg:
                self._place(dph, top=Inches(date_top + shift))
            if aph is not None and ag:
                self._place(aph, top=Inches(author_top + shift))

    def section(self, title, eyebrow=None, description=None, notes=None):
        s = self._new("section")
        m = PH["section"]
        self._refill(s, m["eyebrow"], eyebrow.upper()) if eyebrow else self._drop(s, m["eyebrow"])
        self._refill(s, m["title"], title)
        title_ph = self._ph(s, m["title"])
        if title_ph is not None:
            g = self._layout_geom("section", m["title"])
            if g:
                t_l, t_t, t_w, t_h = g
                lines = min(_est_lines(title, t_w / 914400, 36, True), 3)
                line_h = 0.62
                self._place(title_ph, left=t_l, top=t_t, width=t_w,
                            height=Inches(max(t_h / 914400, lines * line_h)))
                _autofit_shrink(title_ph.text_frame)
                dph = self._ph(s, m["description"])
                dg = self._layout_geom("section", m["description"])
                if description and dph is not None and dg:
                    d_l, d_t, d_w, d_h = dg
                    self._place(dph, left=d_l,
                                top=Inches((t_t / 914400) + lines * line_h + 0.22),
                                width=d_w, height=d_h)
        if description:
            self._refill(s, m["description"], description)
        else:
            self._drop(s, m["description"])
        self.set_notes(s, notes)
        return s

    def content(self, title, body=None, eyebrow=None, lead=None,
                variant="white", notes=None):
        role = "steel" if str(variant).lower().startswith("st") else "white"
        s = self._new(role)
        m = PH["content"]
        self._refill(s, m["title"], title)
        self._refill(s, m["eyebrow"], eyebrow.upper()) if eyebrow else self._drop(s, m["eyebrow"])
        bodyph = self._ph(s, m["body"])
        if bodyph is None:
            return s
        tf = bodyph.text_frame
        tf.word_wrap = True
        # clear example content, keep the placeholder + its native list style
        for p in tf.paragraphs[1:]:
            p._p.getparent().remove(p._p)
        for r in list(tf.paragraphs[0].runs):
            r._r.getparent().remove(r._r)
        items = []
        if lead:
            items.append(("lead", lead))
        for it in (body or []):
            if isinstance(it, (list, tuple)):
                items.append((int(it[0]), str(it[1])))
            else:
                items.append((1, str(it)))
        first = True
        for kind, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = text                       # inherits native bullet list
            if kind == "lead":
                p.level = 0
                _no_bullet(p)
            else:
                p.level = max(0, int(kind) - 1)
        if first:
            self._drop(s, m["body"])
        s._euricom_variant = role
        self.set_notes(s, notes)
        return s

    def content_blank(self, title, eyebrow=None, variant="white", notes=None):
        """Title + eyebrow only; body removed and content area cleared, giving
        a clean CANVAS for a free-drawn component."""
        role = "steel" if str(variant).lower().startswith("st") else "white"
        s = self._new(role)
        m = PH["content"]
        self._refill(s, m["title"], title)
        self._refill(s, m["eyebrow"], eyebrow.upper()) if eyebrow else self._drop(s, m["eyebrow"])
        self._drop(s, m["body"])
        s._euricom_variant = role
        self.set_notes(s, notes)
        return s

    def quote(self, quote, author=None, role=None, notes=None):
        s = self._new("quote")
        m = PH["quote"]
        self._refill(s, m["quote"], quote)
        qph = self._ph(s, m["quote"])
        if qph is not None:
            n = len(quote or "")
            size = 32 if n <= 70 else 26 if n <= 110 else 22 if n <= 160 else 18
            for para in qph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(size)
            _autofit_shrink(qph.text_frame)
        self._refill(s, m["author"], author) if author else self._drop(s, m["author"])
        self._refill(s, m["role"], role) if role else self._drop(s, m["role"])
        self.set_notes(s, notes)
        return s

    # -- speaker notes -----------------------------------------------------
    def set_notes(self, slide, text):
        """Attach speaker notes to a slide's notes pane (the text the presenter
        reads, NOT shown on the slide). Use short cues — 2-4 sentences: the core
        message, what to say out loud, and any figure or transition to land.
        Write notes for every slide; cloned slides start with an empty notes
        pane, so nothing leaks from the template examples. Pass falsy text to
        leave the notes empty."""
        if not text:
            return slide
        slide.notes_slide.notes_text_frame.text = str(text).strip()
        return slide

    # -- component reuse ---------------------------------------------------
    def use_component(self, selector):
        chosen = None
        if isinstance(selector, int):
            chosen = next((c for c in self._catalogue
                           if c["index"] == selector), None)
        else:
            ss = selector.lower()
            chosen = next((c for c in self._catalogue
                           if ss in c["label"].lower() or ss in c["title"].lower()
                           or ss in c["notes"].lower()), None)
        if chosen is None:
            raise ValueError(f"No component matches {selector!r}. Available: "
                             f"{[c['label'] for c in self._catalogue]}")
        s = _duplicate_slide(self.prs, chosen["_slide"])
        self._strip_content_area(s)
        return s

    def set_title(self, slide, title=None, eyebrow=None):
        if title is not None:
            self._refill(slide, PH["content"]["title"], title)
        if eyebrow is not None:
            self._refill(slide, PH["content"]["eyebrow"], eyebrow.upper())
        return slide

    def set_table_rows(self, slide, rows):
        tbl = next((sh.table for sh in slide.shapes if sh.has_table), None)
        if tbl is None:
            raise ValueError("No table on this slide to fill.")
        tref = tbl._tbl
        existing, target = len(tbl.rows), len(rows)
        if target < existing:
            for _ in range(existing - target):
                tref.remove(tref.findall(qn("a:tr"))[-1])
        elif target > existing:
            template_tr = tref.findall(qn("a:tr"))[-1]
            for _ in range(target - existing):
                tref.append(copy.deepcopy(template_tr))
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                if c >= len(tbl.columns):
                    break
                p = tbl.cell(r, c).text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].text = str(val)
                    for extra in p.runs[1:]:
                        extra._r.getparent().remove(extra._r)
                else:
                    p.add_run().text = str(val)
        return slide

    def table(self, slide, headers, rows, x=None, y=None, w=None, h=None,
              total_row=False):
        from components import data_table
        return data_table(slide, headers, rows, x, y, w, h, total_row)

    # -- output ------------------------------------------------------------
    def _cleanup(self):
        lst = self.prs.slides._sldIdLst
        for sid in list(lst):
            if sid.get(qn("r:id")) in self._original_rids:
                try:
                    self.prs.part.drop_rel(sid.get(qn("r:id")))
                except KeyError:
                    pass
                lst.remove(sid)
        root = self.prs.part._element
        for ext in list(root.iter(qn("p:ext"))):
            if ext.find(_P14_SECTIONLST) is not None:
                ext.getparent().remove(ext)

    def save(self, path):
        self._cleanup()
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)
        return path
