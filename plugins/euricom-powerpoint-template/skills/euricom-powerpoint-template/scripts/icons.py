#!/usr/bin/env python3
"""
icons.py - a compact, on-brand monoline icon set for Euricom decks.

Each icon is drawn from native PowerPoint shapes (straight connectors, ovals,
rounded rectangles, a star autoshape) inside a square bbox. That keeps icons
vector, recolourable and editable - no external assets. Default colour is
Midnight; pass the fluo green for an accent.

    import icons
    icons.draw(slide, "rocket", x=1.0, y=2.0, size=0.5)              # Midnight
    icons.draw(slide, "target", x=2.0, y=2.0, size=0.5, color="00FF00")
    icons.NAMES        # list of available icon names

Pick an icon by meaning with icons.suggest("growth") -> "trend".
"""
from __future__ import annotations

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

MIDNIGHT = "014046"


def _stroke(shape, color, w_pt, cap="rnd"):
    ln = shape.line
    ln.color.rgb = RGBColor.from_string(color)
    ln.width = Pt(w_pt)
    lnEl = shape._element.spPr.find(qn("a:ln"))
    if lnEl is not None:
        lnEl.set("cap", cap)
        for tag in ("a:round", "a:bevel", "a:miter"):
            for e in lnEl.findall(qn(tag)):
                lnEl.remove(e)
        lnEl.append(lnEl.makeelement(qn("a:round"), {}))
    if shape.shadow:
        shape.shadow.inherit = False


def _seg(slide, x1, y1, x2, y2, color, w):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _stroke(c, color, w)
    return c


def _oval(slide, x, y, w, h, color, wt, fill=False):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    if fill:
        o.fill.solid(); o.fill.fore_color.rgb = RGBColor.from_string(color)
        o.line.fill.background()
        if o.shadow: o.shadow.inherit = False
    else:
        o.fill.background()
        _stroke(o, color, wt)
    return o


def _rrect(slide, x, y, w, h, color, wt, radius=0.18):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    r.fill.background()
    _stroke(r, color, wt)
    try:
        r.adjustments[0] = radius
    except Exception:
        pass
    return r


# --- individual icons: each fills bbox (x, y, s) ---------------------------
def _check(s, x, y, sz, c, w):
    _seg(s, x + 0.14*sz, y + 0.55*sz, x + 0.40*sz, y + 0.80*sz, c, w)
    _seg(s, x + 0.40*sz, y + 0.80*sz, x + 0.86*sz, y + 0.22*sz, c, w)


def _arrow(s, x, y, sz, c, w):
    _seg(s, x + 0.12*sz, y + 0.5*sz, x + 0.84*sz, y + 0.5*sz, c, w)
    _seg(s, x + 0.58*sz, y + 0.26*sz, x + 0.86*sz, y + 0.5*sz, c, w)
    _seg(s, x + 0.58*sz, y + 0.74*sz, x + 0.86*sz, y + 0.5*sz, c, w)


def _target(s, x, y, sz, c, w):
    _oval(s, x + 0.10*sz, y + 0.10*sz, 0.80*sz, 0.80*sz, c, w)
    _oval(s, x + 0.30*sz, y + 0.30*sz, 0.40*sz, 0.40*sz, c, w)
    _oval(s, x + 0.44*sz, y + 0.44*sz, 0.12*sz, 0.12*sz, c, w, fill=True)


def _bars(s, x, y, sz, c, w):
    base = y + 0.84*sz
    for i, h in enumerate((0.30, 0.52, 0.74)):
        bx = x + 0.18*sz + i*0.26*sz
        _seg(s, bx, base, bx, y + (0.84 - h)*sz, c, w*1.3)


def _trend(s, x, y, sz, c, w):
    _seg(s, x + 0.12*sz, y + 0.74*sz, x + 0.40*sz, y + 0.50*sz, c, w)
    _seg(s, x + 0.40*sz, y + 0.50*sz, x + 0.58*sz, y + 0.60*sz, c, w)
    _seg(s, x + 0.58*sz, y + 0.60*sz, x + 0.86*sz, y + 0.24*sz, c, w)
    _seg(s, x + 0.66*sz, y + 0.24*sz, x + 0.86*sz, y + 0.24*sz, c, w)
    _seg(s, x + 0.86*sz, y + 0.24*sz, x + 0.86*sz, y + 0.44*sz, c, w)


def _clock(s, x, y, sz, c, w):
    _oval(s, x + 0.12*sz, y + 0.12*sz, 0.76*sz, 0.76*sz, c, w)
    _seg(s, x + 0.5*sz, y + 0.5*sz, x + 0.5*sz, y + 0.26*sz, c, w)
    _seg(s, x + 0.5*sz, y + 0.5*sz, x + 0.68*sz, y + 0.58*sz, c, w)


def _layers(s, x, y, sz, c, w):
    for i, dy in enumerate((0.22, 0.44, 0.66)):
        _rrect(s, x + 0.16*sz, y + dy*sz, 0.68*sz, 0.16*sz, c, w, radius=0.3)


def _lock(s, x, y, sz, c, w):
    _rrect(s, x + 0.22*sz, y + 0.44*sz, 0.56*sz, 0.42*sz, c, w, radius=0.18)
    arc = s.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(x + 0.30*sz),
                             Inches(y + 0.14*sz), Inches(0.40*sz), Inches(0.46*sz))
    arc.fill.background(); _stroke(arc, c, w)


def _flag(s, x, y, sz, c, w):
    _seg(s, x + 0.26*sz, y + 0.12*sz, x + 0.26*sz, y + 0.88*sz, c, w)
    _seg(s, x + 0.26*sz, y + 0.16*sz, x + 0.80*sz, y + 0.16*sz, c, w)
    _seg(s, x + 0.80*sz, y + 0.16*sz, x + 0.80*sz, y + 0.46*sz, c, w)
    _seg(s, x + 0.80*sz, y + 0.46*sz, x + 0.26*sz, y + 0.46*sz, c, w)


def _star(s, x, y, sz, c, w):
    st = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x + 0.1*sz),
                            Inches(y + 0.12*sz), Inches(0.80*sz), Inches(0.78*sz))
    st.fill.background(); _stroke(st, c, w)


def _shield(s, x, y, sz, c, w):
    sh = s.shapes.add_shape(MSO_SHAPE.SHIELD if hasattr(MSO_SHAPE, "SHIELD")
                            else MSO_SHAPE.PENTAGON, Inches(x + 0.16*sz),
                            Inches(y + 0.12*sz), Inches(0.68*sz), Inches(0.76*sz))
    sh.fill.background(); _stroke(sh, c, w)
    _check(s, x + 0.14*sz, y + 0.04*sz, 0.72*sz, c, w*0.85)


def _bulb(s, x, y, sz, c, w):
    _oval(s, x + 0.24*sz, y + 0.14*sz, 0.52*sz, 0.52*sz, c, w)
    _seg(s, x + 0.40*sz, y + 0.70*sz, x + 0.60*sz, y + 0.70*sz, c, w)
    _seg(s, x + 0.42*sz, y + 0.80*sz, x + 0.58*sz, y + 0.80*sz, c, w)


def _users(s, x, y, sz, c, w):
    _oval(s, x + 0.30*sz, y + 0.16*sz, 0.30*sz, 0.30*sz, c, w)
    arc = s.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(x + 0.20*sz),
                             Inches(y + 0.52*sz), Inches(0.50*sz), Inches(0.50*sz))
    arc.fill.background(); _stroke(arc, c, w)


def _doc(s, x, y, sz, c, w):
    _rrect(s, x + 0.24*sz, y + 0.12*sz, 0.52*sz, 0.76*sz, c, w, radius=0.06)
    for dy in (0.36, 0.50, 0.64):
        _seg(s, x + 0.34*sz, y + dy*sz, x + 0.66*sz, y + dy*sz, c, w*0.8)


def _gear(s, x, y, sz, c, w):
    _oval(s, x + 0.24*sz, y + 0.24*sz, 0.52*sz, 0.52*sz, c, w)
    _oval(s, x + 0.40*sz, y + 0.40*sz, 0.20*sz, 0.20*sz, c, w)
    import math
    cx, cy, r1, r2 = x + 0.5*sz, y + 0.5*sz, 0.26*sz, 0.36*sz
    for k in range(8):
        a = k * math.pi / 4
        _seg(s, cx + r1*math.cos(a), cy + r1*math.sin(a),
             cx + r2*math.cos(a), cy + r2*math.sin(a), c, w)


_ICONS = {
    "check": _check, "arrow": _arrow, "target": _target, "bars": _bars,
    "trend": _trend, "clock": _clock, "layers": _layers, "lock": _lock,
    "flag": _flag, "star": _star, "shield": _shield, "bulb": _bulb,
    "users": _users, "doc": _doc, "gear": _gear,
}
NAMES = sorted(_ICONS)

# meaning -> icon (rough; extend freely)
_SYNONYMS = {
    "growth": "trend", "increase": "trend", "revenue": "trend", "kpi": "bars",
    "data": "bars", "metric": "bars", "goal": "target", "focus": "target",
    "time": "clock", "speed": "trend", "deadline": "clock", "stack": "layers",
    "platform": "layers", "architecture": "layers", "security": "shield",
    "safe": "shield", "privacy": "lock", "access": "lock", "milestone": "flag",
    "quality": "star", "premium": "star", "idea": "bulb", "innovation": "bulb",
    "team": "users", "customer": "users", "people": "users", "document": "doc",
    "report": "doc", "process": "gear", "automation": "gear", "config": "gear",
    "done": "check", "success": "check", "next": "arrow", "step": "arrow",
}


def suggest(keyword):
    k = (keyword or "").lower().strip()
    if k in _ICONS:
        return k
    return _SYNONYMS.get(k, "check")


def draw(slide, name, x, y, size=0.5, color=MIDNIGHT, weight=None):
    fn = _ICONS.get(name) or _ICONS[suggest(name)]
    w = weight if weight is not None else max(1.4, size * 72 * 0.055)
    fn(slide, x, y, size, color, w)
